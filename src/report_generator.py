"""
report_generator.py — Resident report PPTX builder
====================================================

Generates one PPTX resident report per site_id from the master XRF CSV
and site_configs.json, by filling placeholders in
``src/Resident_Report_Template.pptx``.

This module is intentionally framework-free (no streamlit imports) so
it can be:
  * called from ``etl_manager.py`` (the Streamlit pipeline UI)
  * called from any standalone script or unit test
  * driven from a CLI / cron job

Template slide layout (1-indexed page numbers as the user sees them in
PowerPoint; 0-indexed when iterating with python-pptx):

  Page 1 / slide_idx 0 — Cover letter
      * "Address of Resident" → SiteID
      * "Name of Resident"    → SiteID
      * "Date"                → today
  Page 2 / slide_idx 1 — Sample collection method (no edits)
  Page 3 / slide_idx 2 — Soil Report Summary
      * Inserts a basemap-style map (Leaflet HTML screenshot with
        numbers shown). Falls back to the light-theme static PNG if
        the Leaflet screenshot can't be produced (e.g. no internet for
        tile servers).
      * Strips the "Heat map of property (no basemap)" placeholder.
  Page 4 / slide_idx 3 — Detailed Results
      * Backyard / Frontyard PPM values filled into the existing table.
      * If a yard has no data, the displayed value is "-".
      * Inserts the **white-background** static PNG (light_theme=True).
      * Strips the "Heat map of property with basemap" placeholder.
  Page 5+ — keep as-is.

Data policy
-----------
Only real XRF readings produce numbers in reports.  mock_ppm values
from site_configs.json are NEVER written into a generated report.
A site with zero real samples is skipped entirely (we don't want
placeholder numbers in something a resident reads).
"""

from __future__ import annotations

import json
import math
import os
import sys
from datetime import date
from typing import Optional, Tuple, List, Dict

import pandas as pd
from pptx import Presentation
from pptx.util import Inches

# Make sibling modules importable regardless of where this is called from.
_HERE = os.path.dirname(os.path.abspath(__file__))
if _HERE not in sys.path:
    sys.path.insert(0, _HERE)

from groundsense_config import get_nysh_category, resolve_lod
from map_renderer import render_static_png, render_leaflet_html


# ───────────────────────────────────────────────────────────
#  ZONE / PPM HELPERS
# ───────────────────────────────────────────────────────────
def _flatten_grid_blocks(site_config: dict) -> Dict[str, dict]:
    """Return a flat ``{block_id: block_dict}`` map regardless of whether
    the site uses the legacy top-level ``grid_blocks`` layout or the
    newer ``yards.<yk>.grid_blocks`` layout.

    For yard-style configs we propagate the yard key down onto each
    block as both ``yard`` and ``zone`` (if missing) so downstream code
    doesn't need to know about the layout.
    """
    yards = site_config.get("yards") or {}
    if yards:
        flat: Dict[str, dict] = {}
        for yk, y in yards.items():
            for bid, b in (y.get("grid_blocks") or {}).items():
                b = dict(b)
                b.setdefault("yard", yk)
                b.setdefault("zone", yk)
                flat[bid] = b
        if flat:
            return flat
    return dict(site_config.get("grid_blocks", {}))


def _block_real_ppm(block: dict, master_df: pd.DataFrame) -> Optional[float]:
    """Average LeadPPM_Clean for a single block, or None if no real
    sample IDs matched.  Uses the same pattern-matching strategy as
    map_renderer.
    """
    patterns = block.get("sample_id_patterns", []) or []
    for pat in patterns:
        if not pat:
            continue
        matches = master_df[
            master_df['SampleID'].str.contains(pat, case=False, na=False)
        ]
        if not matches.empty:
            avg = matches['LeadPPM_Clean'].mean()
            if pd.notna(avg):
                return float(avg)
    return None


def compute_zone_averages(site_config: dict, master_df: pd.DataFrame) -> dict:
    """Average LeadPPM per zone using only real XRF data.

    Returns e.g. ``{"back": 542.3, "front": 718.1}``.  A zone is absent
    from the result dict if no SampleID pattern in that zone matched
    a real reading.

    Kept for backwards compatibility with the previous etl_manager API.
    New code should prefer ``compute_zone_stats`` (which also exposes
    min/max per yard for the page-4 table).
    """
    if 'LeadPPM_Clean' not in master_df.columns:
        master_df = master_df.copy()
        master_df['LeadPPM_Clean'] = master_df['LeadPPM'].apply(resolve_lod)

    grid = _flatten_grid_blocks(site_config)
    zone_values: Dict[str, List[float]] = {}

    for block_id, dims in grid.items():
        if block_id.startswith("_"):
            continue
        zone = dims.get("zone") or dims.get("yard") or "yard"
        ppm = _block_real_ppm(dims, master_df)
        if ppm is not None and not math.isnan(ppm):
            zone_values.setdefault(zone, []).append(ppm)

    return {z: sum(v) / len(v) for z, v in zone_values.items() if v}


def compute_zone_stats(
    site_config: dict,
    master_df: pd.DataFrame,
) -> Dict[str, Dict[str, Optional[float]]]:
    """Per-yard {highest, lowest, average} of real block-level PPM averages.

    Each block in the site has its own average PPM (mean of all readings
    matching its SampleID patterns).  The page-4 table in the resident
    report wants the *spread* of those block averages within each yard
    — i.e. across the site's grid cells, what's the highest reading, the
    lowest, and the mean.  Blocks with no real data are excluded.

    Returns
    -------
    dict
        Keyed by canonical yard name (``"backyard"`` / ``"frontyard"``).
        Each value is ``{"highest": float|None, "lowest": float|None,
        "average": float|None, "count": int}``.  A yard with no real
        readings has all-None stats and count=0.
    """
    if 'LeadPPM_Clean' not in master_df.columns:
        master_df = master_df.copy()
        master_df['LeadPPM_Clean'] = master_df['LeadPPM'].apply(resolve_lod)

    grid = _flatten_grid_blocks(site_config)

    # Bucket per-block real PPMs by canonical yard.
    by_yard: Dict[str, List[float]] = {"backyard": [], "frontyard": []}
    for block_id, dims in grid.items():
        if block_id.startswith("_"):
            continue
        zone = (dims.get("zone") or dims.get("yard") or "yard").lower()
        canonical = _canonical_yard(zone)
        if canonical is None:
            continue
        ppm = _block_real_ppm(dims, master_df)
        if ppm is not None and not math.isnan(ppm):
            by_yard[canonical].append(ppm)

    def _stats(vals: List[float]) -> Dict[str, Optional[float]]:
        if not vals:
            return {"highest": None, "lowest": None,
                    "average": None, "count": 0}
        return {
            "highest": max(vals),
            "lowest": min(vals),
            "average": sum(vals) / len(vals),
            "count": len(vals),
        }

    return {y: _stats(vs) for y, vs in by_yard.items()}


def _canonical_yard(zone_name: str) -> Optional[str]:
    """Map any zone/yard label used in site_configs to ``"backyard"``,
    ``"frontyard"``, or None (skip).

    Single-zone sites (``yard``, ``transect``) collapse into the
    backyard bucket since those reports have no front-yard row data
    to populate — the front-yard cells will end up as ``"-"``.
    """
    z = (zone_name or "").lower()
    if z in ("back", "backyard", "back_yard"):
        return "backyard"
    if z in ("front", "frontyard", "front_yard"):
        return "frontyard"
    if z in ("yard", "transect"):
        return "backyard"
    return None


def format_zone_ppm(zone_averages: dict) -> Dict[str, Optional[int]]:
    """Map raw zone averages → {backyard_ppm, frontyard_ppm} (ints or None).

    Single-zone sites (``yard``, ``transect``) collapse into the
    backyard slot since those reports only have one map / one number.

    Retained for backwards compatibility; the new page-4 table uses
    ``compute_zone_stats`` directly.
    """
    result: Dict[str, Optional[int]] = {
        "backyard_ppm": None,
        "frontyard_ppm": None,
    }
    for zone, avg in zone_averages.items():
        canonical = _canonical_yard(zone)
        if canonical == "backyard" and result["backyard_ppm"] is None:
            result["backyard_ppm"] = round(avg)
        elif canonical == "frontyard":
            result["frontyard_ppm"] = round(avg)
    return result


# ───────────────────────────────────────────────────────────
#  MAP IMAGE RENDERING
# ───────────────────────────────────────────────────────────
# CSS injected into the Leaflet HTML before screenshotting so that
# the interactive panels (which the resident doesn't need to see) are
# hidden in the printed report.  The map + legend remain visible.
_REPORT_SCREENSHOT_CSS = """
.title-bar, .controls, .anchor-toggle { display: none !important; }
body { background: #ffffff !important; }
#map { background: #ffffff !important; }
"""


def render_basemap_screenshot(
    site_config: dict,
    master_df: pd.DataFrame,
    output_path: str,
    timeout_ms: int = 15000,
    wait_for_tiles_ms: int = 6000,
) -> bool:
    """Render the Leaflet basemap+numbers HTML as a static PNG.

    Used for slide 3 of the report ("Map of Site" — basemap + numbers).
    Requires playwright + a chromium browser, and live internet so
    OpenStreetMap/satellite tiles can load.

    Returns
    -------
    bool
        True on success, False if anything went wrong (caller should
        fall back to the static no-basemap PNG).
    """
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        return False

    try:
        html = render_leaflet_html(
            site_config, master_df,
            show_numbers=True,         # page 3 wants the values visible
            use_mock_fallback=False,   # real data only in reports
        )
    except Exception:
        return False

    # Inject our report-mode CSS so the floating control panel and
    # title bar don't appear in the screenshot.  We append a <style>
    # block right before </head>; if for some reason there's no </head>
    # we just prepend the style to the body.
    style_tag = f"<style>{_REPORT_SCREENSHOT_CSS}</style>"
    if "</head>" in html:
        html = html.replace("</head>", style_tag + "</head>", 1)
    else:
        html = style_tag + html

    # Write to a temp file next to the output so file:// loads correctly
    # (some browsers refuse data: URLs the size of a full Leaflet doc).
    out_dir = os.path.dirname(os.path.abspath(output_path)) or "."
    os.makedirs(out_dir, exist_ok=True)
    tmp_html = os.path.join(out_dir, f".{os.path.basename(output_path)}.html")
    with open(tmp_html, "w", encoding="utf-8") as f:
        f.write(html)

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            try:
                ctx = browser.new_context(viewport={"width": 1200, "height": 900})
                page = ctx.new_page()
                page.goto(
                    "file://" + os.path.abspath(tmp_html),
                    wait_until="domcontentloaded",
                    timeout=timeout_ms,
                )
                # Give Leaflet time to lay out + tiles to load.  We don't
                # block on networkidle because some tile servers stream
                # forever; a fixed wait is more predictable.
                page.wait_for_timeout(wait_for_tiles_ms)
                page.screenshot(path=output_path, full_page=False)
            finally:
                browser.close()
    except Exception:
        # Any browser/network failure → caller falls back to static PNG.
        try:
            os.remove(tmp_html)
        except OSError:
            pass
        return False

    # Heuristic: if Playwright couldn't fetch tiles (no internet) the
    # map area will be mostly one solid colour.  We sample the *inner*
    # central region — small enough that the legend (bottom-left) and
    # any chrome (top corners) is fully outside the crop — and count
    # distinct coarse colour buckets.  A real basemap with site cells
    # produces many buckets (dozens); an empty Leaflet canvas produces
    # 1–3.  Threshold of 6 leaves comfortable headroom on both sides.
    try:
        from PIL import Image
        with Image.open(output_path) as im:
            w, h = im.size
            box = (int(w * 0.30), int(h * 0.25),
                   int(w * 0.70), int(h * 0.65))
            cropped = im.convert("RGB").crop(box).resize((80, 60))
            pixels = list(cropped.getdata())
            buckets = {(r >> 5, g >> 5, b >> 5) for r, g, b in pixels}
            if len(buckets) < 6:
                # Mostly featureless — tiles never loaded.
                try:
                    os.remove(tmp_html)
                except OSError:
                    pass
                return False
    except Exception:
        # PIL missing or unreadable image — assume the screenshot is OK.
        pass

    try:
        os.remove(tmp_html)
    except OSError:
        pass
    return True


def render_white_static_map(
    site_config: dict,
    master_df: pd.DataFrame,
    output_path: str,
) -> bool:
    """Render the **white-background** static PNG (light_theme=True).

    Used for slide 4 of the report ("Detailed Results" — heat map with
    no basemap).  Also used as the fallback for slide 3 when the
    Leaflet screenshot can't be produced.

    Returns True on success.
    """
    try:
        render_static_png(
            site_config, master_df, output_path,
            show_numbers=True,
            use_mock_fallback=False,
            light_theme=True,
        )
    except Exception:
        return False
    return os.path.exists(output_path)


# ───────────────────────────────────────────────────────────
#  TEXT PLACEHOLDER REWRITING
# ───────────────────────────────────────────────────────────
# Phrases that exist purely as designer placeholders in the template
# and should be stripped from any generated report.  These are matched
# at the *paragraph* level (joined run text), not per-run, so that
# notes split across multiple runs are still caught.
_PLACEHOLDER_PARAGRAPHS_TO_STRIP = (
    # Current template (uploaded May 2026)
    "Heat map of property with basemap and numbers",
    "Heat map of property (no basemap)",
    # Legacy phrases retained for backwards compatibility with older
    # versions of Resident_Report_Template.pptx.
    "Heat map of property with basemap",
    "Visual map of property with color-coded zones",
    "Highlight hotspots",
)

# Placeholder strings on page 1 that should be replaced with the SiteID.
# Order matters: longer / more specific strings come first so that
# "Address of Resident from Site Analysis Database" is matched before
# the legacy "Address of Resident" prefix that's a substring of it.
_SITE_ID_PLACEHOLDERS = (
    "Address of Resident from Site Analysis Database",
    "Address of Resident",
    "Name of Resident",
)


def _fmt_ppm(value: Optional[float]) -> str:
    """Render a PPM value for the slide-4 table.

    Returns ``"-"`` when no real data was available (so the front-yard
    row of a single-yard site renders as a literal dash, as the user
    requested).  Otherwise emits a plain integer — no decimals, no
    comma separators — matching the template's "### ppm" typography.
    """
    if value is None:
        return "-"
    return str(int(round(value)))


def _strip_placeholder_paragraph(paragraph) -> bool:
    """If the paragraph's combined run text contains a known designer
    placeholder phrase, blank out the entire paragraph's text and
    return True.

    We blank out the paragraph rather than the matching substring
    alone because these placeholders are full multi-line designer
    notes (e.g. "Heat map of property with basemap and numbers – drag
    and crop image to ensure that front street, backyard, and key are
    all visible") that span several runs, and leaving partial text
    behind would look broken.
    """
    para_text = "".join(r.text or "" for r in paragraph.runs)
    for phrase in _PLACEHOLDER_PARAGRAPHS_TO_STRIP:
        if phrase in para_text:
            for run in paragraph.runs:
                if run.text:
                    run.text = ""
            return True
    return False


def _replace_in_run(run, old: str, new: str) -> None:
    """Replace ``old`` with ``new`` inside a single run's text,
    preserving the run's formatting.
    """
    if run.text and old in run.text:
        run.text = run.text.replace(old, new)


def _rewrite_text_frame(
    text_frame,
    slide_idx: int,
    site_id: str,
    today_str: str,
    table_value_queue: List[str],
) -> None:
    """Apply all per-text-frame substitutions.

    Used for both top-level shapes and table cells (a table cell has
    its own text_frame).  The slide-4 ``###`` substitution is driven
    by a positional queue prepared by the caller — every ``###`` token
    encountered (left-to-right, top-to-bottom) consumes the next value
    off the queue.  This keeps the substitution order deterministic
    even though the runs themselves don't carry any "which cell?"
    metadata.
    """
    for paragraph in text_frame.paragraphs:
        # Strip multi-run designer-note placeholders first — those are
        # the "Heat map of property…" reminders that should never appear
        # in a generated report.
        if _strip_placeholder_paragraph(paragraph):
            continue

        for run in paragraph.runs:
            txt = run.text
            if not txt:
                continue

            # ── Slide 1 (cover letter) ─────────────────────────────
            if slide_idx == 0:
                for placeholder in _SITE_ID_PLACEHOLDERS:
                    if placeholder in txt:
                        txt = txt.replace(placeholder, site_id)
                if txt.strip() == "Date":
                    txt = today_str

            # ── Slide 4 table cells: replace ### with next queue value
            # Only triggered on slide 4 (slide_idx==3) where the caller
            # primed the queue.  Multiple ### in one run are handled by
            # the while loop.
            if slide_idx == 3:
                while "###" in txt and table_value_queue:
                    txt = txt.replace("###", table_value_queue.pop(0), 1)

            if txt != run.text:
                run.text = txt


def _rewrite_slide(
    slide,
    slide_idx: int,
    site_id: str,
    today_str: str,
    zone_stats: Dict[str, Dict[str, Optional[float]]],
) -> None:
    """Walk every shape on the slide and apply the right substitutions.

    For slide 4 specifically, we recognise the
    Backyard/Front yard × Highest/Lowest/Average table and fill its
    six data cells.  We do this by:

      1. Locating the table inside the slide's shapes (python-pptx
         exposes it via ``shape.has_table`` / ``shape.table``).
      2. Reading each row's leading cell text to identify whether the
         row is the Backyard row or the Front yard row.
      3. Assigning {highest, lowest, average} values to the remaining
         cells *in their visual order* (columns 1..3 of that row).

    This approach doesn't rely on column-header text matching, which
    avoids fragility if the template's header row ever gets reformatted
    or split across runs.  It just trusts the column order
    (Highest, Lowest, Average) which is fixed by the template design.

    Outside the table we just rewrite text frames the normal way.
    """
    backyard = zone_stats.get("backyard", {})
    frontyard = zone_stats.get("frontyard", {})

    # Build a positional queue of replacement values for the 6 ###
    # tokens we expect to find inside the table, in the order
    # python-pptx walks rows × columns.  The template has 3 rows
    # (header + Backyard + Front yard) and 4 columns ((row label) +
    # Highest + Lowest + Average).  The header row has no ### tokens,
    # so the queue maps onto the 6 data cells (Backyard's H/L/A then
    # Front yard's H/L/A).
    yard_value_queue: List[str] = [
        _fmt_ppm(backyard.get("highest")),
        _fmt_ppm(backyard.get("lowest")),
        _fmt_ppm(backyard.get("average")),
        _fmt_ppm(frontyard.get("highest")),
        _fmt_ppm(frontyard.get("lowest")),
        _fmt_ppm(frontyard.get("average")),
    ]

    for shape in slide.shapes:
        # ── Table shapes (the page-4 results table) ───────────────
        if getattr(shape, "has_table", False) and shape.has_table:
            # Skip the header row entirely — no ### tokens to fill,
            # and we don't want to accidentally treat row-label cells
            # as data cells.
            for row_idx, row in enumerate(shape.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    if cell.text_frame is None:
                        continue
                    # Row-label cells (column 0 on Backyard/Front yard
                    # rows) sometimes contain no ### token, so the
                    # queue stays untouched.  Data cells consume from
                    # the queue inside _rewrite_text_frame.
                    _rewrite_text_frame(
                        cell.text_frame,
                        slide_idx=slide_idx,
                        site_id=site_id,
                        today_str=today_str,
                        table_value_queue=yard_value_queue,
                    )
            continue

        # ── Regular text frame shapes ─────────────────────────────
        if not shape.has_text_frame:
            continue
        _rewrite_text_frame(
            shape.text_frame,
            slide_idx=slide_idx,
            site_id=site_id,
            today_str=today_str,
            table_value_queue=yard_value_queue,
        )


# ───────────────────────────────────────────────────────────
#  MAP IMAGE INSERTION
# ───────────────────────────────────────────────────────────
def _insert_picture_centered(
    slide,
    image_path: str,
    *,
    slide_width_in: float,
    max_w_in: float,
    max_h_in: float,
    top_in: float,
) -> None:
    """Add a picture to the slide, centered horizontally inside a
    bounding box of (max_w_in × max_h_in), preserving aspect ratio.
    """
    try:
        from PIL import Image as PILImage
        with PILImage.open(image_path) as img:
            img_aspect = img.width / img.height
    except Exception:
        img_aspect = 10.0 / 7.5  # safe default for matplotlib output

    if max_w_in / img_aspect > max_h_in:
        # Height-limited
        h = max_h_in
        w = max_h_in * img_aspect
    else:
        # Width-limited
        w = max_w_in
        h = max_w_in / img_aspect

    left = Inches((slide_width_in - w) / 2)
    slide.shapes.add_picture(
        image_path, left, Inches(top_in),
        width=Inches(w), height=Inches(h),
    )


# ───────────────────────────────────────────────────────────
#  PUBLIC API
# ───────────────────────────────────────────────────────────
def generate_report_for_site(
    site_config: dict,
    master_df: pd.DataFrame,
    template_path: str,
    output_dir: str,
    maps_dir: Optional[str] = None,
    today_str: Optional[str] = None,
    log=print,
) -> Optional[Tuple[str, str]]:
    """Generate one resident report for a single site.

    Returns ``(site_id, output_pptx_path)`` on success, or ``None`` if
    the site was skipped (no real data) or generation failed.
    """
    site_id = site_config.get("site_id", "UnknownSite")
    today_str = today_str or date.today().strftime("%m/%d/%Y")

    if maps_dir is None:
        maps_dir = os.path.join(output_dir, "map_images")
    os.makedirs(output_dir, exist_ok=True)
    os.makedirs(maps_dir, exist_ok=True)

    # ── 1. Compute PPM stats from real data only ────────────────
    # zone_stats is the source of truth for page-4 (per-yard
    # Highest/Lowest/Average).  We also keep zone_averages around so
    # the "skip if no real data" decision and the "any real number?"
    # check below remain identical to the legacy implementation.
    zone_stats = compute_zone_stats(site_config, master_df)
    zone_averages = compute_zone_averages(site_config, master_df)
    if not zone_averages:
        log(f"⏭️  Skipped '{site_id}' — no real XRF data matched yet.")
        return None

    # ── 2. Render map images ────────────────────────────────────
    # Page 4 — white static heat map (always rendered, used as both
    # the page-4 map and the page-3 fallback).
    static_map_path = os.path.join(maps_dir, f"map_{site_id}_static.png")
    static_ok = render_white_static_map(site_config, master_df, static_map_path)
    if not static_ok:
        log(f"⚠️  Could not render static map for {site_id}; report will lack map images.")
        static_map_path = None

    # Page 3 — Leaflet basemap screenshot (with numbers).  Falls back
    # to the static white map if the screenshot can't be produced.
    basemap_map_path = os.path.join(maps_dir, f"map_{site_id}_basemap.png")
    basemap_ok = render_basemap_screenshot(
        site_config, master_df, basemap_map_path
    )
    if not basemap_ok:
        # Use the static map as a fallback so page 3 isn't empty.
        basemap_map_path = static_map_path
        log(
            f"ℹ️  Basemap screenshot unavailable for {site_id} "
            f"(no browser / no tiles); using static map on page 3 instead."
        )

    # ── 3. Open template ────────────────────────────────────────
    try:
        prs = Presentation(template_path)
    except Exception as e:
        log(f"❌ Failed to open template for {site_id}: {e}")
        return None

    # Reference slide width for centering math.  Templates can use
    # widescreen (13.333") or 4:3 (10").  python-pptx returns EMUs.
    slide_width_in = prs.slide_width / 914400.0

    # ── 4. Rewrite text + insert maps ───────────────────────────
    for slide_idx, slide in enumerate(prs.slides):
        _rewrite_slide(
            slide,
            slide_idx=slide_idx,
            site_id=site_id,
            today_str=today_str,
            zone_stats=zone_stats,
        )

        # Page 3 (slide_idx 2) — basemap-with-numbers map.  The template
        # reserves a large block below the "Map of Site" blue banner
        # (banner bottom ≈ y=4.2") but the QR-code "Scan for resources"
        # panel sits at y≈8.8" in the lower-right corner — so the map
        # must stop short of it.  4.2" of height (top 4.35" → bottom
        # ≈ 8.55") leaves a clear margin above the QR panel.
        if slide_idx == 2 and basemap_map_path and os.path.exists(basemap_map_path):
            try:
                _insert_picture_centered(
                    slide, basemap_map_path,
                    slide_width_in=slide_width_in,
                    max_w_in=6.3, max_h_in=4.2,
                    top_in=4.35,
                )
            except Exception as e:
                log(f"⚠️  Map insert failed (page 3) for {site_id}: {e}")

        # Page 4 (slide_idx 3) — white static heat map.  The Detailed
        # Results table ends around y=3.05"; we put the map below it.
        # Slide is 8.5" × 11" portrait so there's ~6" of usable height
        # before the footer.
        if slide_idx == 3 and static_map_path and os.path.exists(static_map_path):
            try:
                _insert_picture_centered(
                    slide, static_map_path,
                    slide_width_in=slide_width_in,
                    max_w_in=6.5, max_h_in=5.4,
                    top_in=3.35,
                )
            except Exception as e:
                log(f"⚠️  Map insert failed (page 4) for {site_id}: {e}")

    # ── 5. Save ────────────────────────────────────────────────
    output_file = os.path.join(output_dir, f"Resident_Report_{site_id}.pptx")
    try:
        prs.save(output_file)
    except Exception as e:
        log(f"❌ Failed to save report for {site_id}: {e}")
        return None
    return site_id, output_file


def generate_all_reports(
    master_csv_path: str,
    template_path: str,
    output_dir: str,
    site_configs_path: str,
    log=print,
) -> Tuple[int, List[Tuple[str, str]]]:
    """Generate one PPTX per site listed in ``site_configs.json``.

    Returns ``(count, [(site_id, path), ...])``.

    Sites with no matching real XRF data are skipped (not a failure).
    """
    df = pd.read_csv(master_csv_path)
    df = df[df['SampleID'].notna() & (df['SampleID'] != "")]
    df['LeadPPM_Clean'] = df['LeadPPM'].apply(resolve_lod)

    if not os.path.exists(site_configs_path):
        log(f"❌ Site configs not found at {site_configs_path}")
        return 0, []

    with open(site_configs_path) as f:
        raw = json.load(f)
    site_configs = raw if isinstance(raw, list) else list(raw.values())

    if not site_configs:
        log("⚠️ No site configurations found — nothing to do.")
        return 0, []

    today_str = date.today().strftime("%m/%d/%Y")
    results: List[Tuple[str, str]] = []

    for site_config in site_configs:
        outcome = generate_report_for_site(
            site_config, df,
            template_path=template_path,
            output_dir=output_dir,
            today_str=today_str,
            log=log,
        )
        if outcome is not None:
            results.append(outcome)
            log(f"✅ Generated report for {outcome[0]}")

    return len(results), results


# ───────────────────────────────────────────────────────────
#  CLI entry point
# ───────────────────────────────────────────────────────────
def _main() -> int:
    import argparse
    ap = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    ap.add_argument("--master-csv", required=True,
                    help="Path to merged master XRF chemistry CSV")
    ap.add_argument("--template",
                    default=os.path.join(_HERE, "Resident_Report_Template.pptx"))
    ap.add_argument("--site-configs", required=True,
                    help="Path to site_configs.json")
    ap.add_argument("--out-dir", required=True,
                    help="Directory to write Resident_Report_*.pptx files into")
    args = ap.parse_args()

    count, results = generate_all_reports(
        master_csv_path=args.master_csv,
        template_path=args.template,
        output_dir=args.out_dir,
        site_configs_path=args.site_configs,
    )
    print(f"\nGenerated {count} report(s).")
    for sid, path in results:
        print(f"  • {sid:>20s}  →  {path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(_main())