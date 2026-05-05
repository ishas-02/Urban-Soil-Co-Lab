
#  """
# etl_manager.py — GroundSense Data Pipeline Manager

# Updates from original:
#   1. Imports NYSH colors & rotation helpers from groundsense_config.py
#   2. **NEW**: Delegates map image generation to `map_renderer.render_static_png()`
#      — the same function site_builder.py uses for its PNG export. PPTX maps
#      are now byte-identical to the "No basemap · with numbers" download.
#   3. Honors `rotation_deg` saved in site_configs.json (via map_renderer)
#   4. Fills in zone-based average Lead PPM values (backyard / front yard)
#   5. Data policy: only real XRF data is used in exported maps. Cells without
#      matching Master Data render as gray "No Data" — no silent mock fallback.

# Template slide layout (6 pages, 0-indexed):
#   Slide 0: Cover letter — replace "Address of Resident", "Name of Resident", "Date"
#   Slide 1: Sample Collection Method — keep as-is
#   Slide 2: Soil Report Summary — insert dark map (no basemap)
#   Slide 3: Detailed Results — fill [###] with real PPM, insert map
#   Slide 4: NY Soil Health info — keep as-is
#   Slide 5: Lead level table & safety — keep as-is
# """


# import streamlit as st
# import os
# import sys
# import glob
# import re
# import subprocess
# import json
# import math
# import pandas as pd
# from pptx import Presentation
# from pptx.util import Inches
# from datetime import date

# # Inject both the script's own directory AND src/ into sys.path so that
# # `groundsense_config` and `map_renderer` can be imported regardless of
# # whether they live in src/ or at the repo root next to this script.
# _HERE = os.path.dirname(os.path.abspath(__file__))
# for _p in (_HERE, os.path.join(_HERE, "src")):
#     if _p not in sys.path:
#         sys.path.insert(0, _p)

# from groundsense_config import (
#     get_nysh_category,
#     NYSH_TIERS,
#     NYSH_COLORS,
#     calculate_coordinate,
#     resolve_lod,
# )
# from map_renderer import render_static_png  # single source of truth for visuals

# # ═══════════════════════════════════════════════
# #  PAGE CONFIGURATION
# # ═══════════════════════════════════════════════
# st.set_page_config(page_title="GroundSense Pipeline", page_icon="⚙️", layout="wide")

# # Initialize Session State
# if 'pipeline_success' not in st.session_state:
#     st.session_state.pipeline_success = False
# if 'latest_master_file' not in st.session_state:
#     st.session_state.latest_master_file = None
# if 'reports_generated' not in st.session_state:
#     st.session_state.reports_generated = 0
# if 'generated_report_list' not in st.session_state:
#     st.session_state.generated_report_list = []


# # ═══════════════════════════════════════════════
# #  ZONE-BASED PPM CALCULATOR
# # ═══════════════════════════════════════════════
# def compute_zone_averages(site_config, master_df):
#     """Compute average Lead PPM for each zone (back, front, yard, transect).

#     Returns dict, e.g. {"back": 542.3, "front": 718.1}

#     IMPORTANT: Only uses real XRF data. mock_ppm values are ignored here
#     because the resident report must show truthful measurements — placeholders
#     belong only in the design-time site_builder preview.
#     """
#     grid = site_config.get("grid_blocks", {})

#     if 'LeadPPM_Clean' not in master_df.columns:
#         master_df = master_df.copy()
#         master_df['LeadPPM_Clean'] = master_df['LeadPPM'].apply(resolve_lod)

#     zone_values = {}  # zone_name -> [ppm, ...]

#     for block_id, dims in grid.items():
#         if block_id.startswith("_"):
#             continue
#         zone = dims.get("zone", "yard")
#         patterns = dims.get("sample_id_patterns", [])

#         ppm = None
#         for pat in patterns:
#             if not pat:
#                 continue
#             matches = master_df[
#                 master_df['SampleID'].str.contains(pat, case=False, na=False)
#             ]
#             if not matches.empty:
#                 avg = matches['LeadPPM_Clean'].mean()
#                 if pd.notna(avg):
#                     ppm = avg
#                     break

#         # Skip mock fallback — real data only
#         if ppm is not None and not (isinstance(ppm, float) and math.isnan(ppm)):
#             zone_values.setdefault(zone, []).append(ppm)

#     return {z: sum(v) / len(v) for z, v in zone_values.items() if v}


# def format_zone_ppm(zone_averages):
#     """Map zone averages to backyard_ppm / frontyard_ppm for template filling.

#     Single-zone sites (yard, transect) map to backyard only.
#     """
#     result = {"backyard_ppm": None, "frontyard_ppm": None}
#     for zone, avg in zone_averages.items():
#         z = zone.lower()
#         if z in ("back", "backyard"):
#             result["backyard_ppm"] = round(avg)
#         elif z in ("front", "frontyard", "front_yard"):
#             result["frontyard_ppm"] = round(avg)
#         elif z in ("yard", "transect"):
#             result["backyard_ppm"] = round(avg)
#     return result


# # ═══════════════════════════════════════════════
# #  REPORT GENERATION (PPTX)
# # ═══════════════════════════════════════════════
# def generate_pptx_reports(master_csv_path, template_path, output_dir,
#                           site_db_path, site_configs_path):
#     """Generate one PPTX resident report per site address.

#     Uses the shared `map_renderer.render_static_png()` so the map image
#     embedded in the PPTX is byte-identical to the PNG download from
#     site_builder.py's "No basemap · with numbers" option.

#     Returns (report_count, list_of (site_address, filepath) tuples).
#     """
#     df = pd.read_csv(master_csv_path)
#     df = df[df['SampleID'].notna() & (df['SampleID'] != "")]
#     df['LeadPPM_Clean'] = df['LeadPPM'].apply(resolve_lod)

#     # Load site configs — one report per site address
#     site_configs = {}
#     if os.path.exists(site_configs_path):
#         with open(site_configs_path, 'r') as f:
#             raw = json.load(f)
#         site_configs = {s["address"]: s for s in raw}

#     if not site_configs:
#         st.warning("⚠️ No site configurations found. Cannot generate reports.")
#         return 0, []

#     os.makedirs(output_dir, exist_ok=True)
#     maps_dir = os.path.join(output_dir, "map_images")
#     os.makedirs(maps_dir, exist_ok=True)

#     report_count = 0
#     generated_reports = []
#     today_str = date.today().strftime("%m/%d/%Y")

#     for site_address, site_config in site_configs.items():

#         # ── 1. Compute zone averages (real data only) ──
#         zone_averages = compute_zone_averages(site_config, df)
#         zone_ppm = format_zone_ppm(zone_averages)

#         all_zone_vals = list(zone_averages.values())
#         if all_zone_vals:
#             site_avg = sum(all_zone_vals) / len(all_zone_vals)
#         else:
#             # No real data for this site — skip report generation entirely.
#             # Previously we fell back to mock_ppm, but that leaks placeholder
#             # numbers into official reports, which we explicitly don't want.
#             st.info(
#                 f"⏭️  Skipped '{site_address}' — no real XRF data matched yet."
#             )
#             continue

#         if zone_ppm["backyard_ppm"] is None:
#             zone_ppm["backyard_ppm"] = round(site_avg)

#         # ── 2. Render static map via shared renderer ──
#         safe_name = "".join(
#             c for c in site_address if c.isalnum() or c == ' '
#         ).rstrip()
#         map_image_path = os.path.join(maps_dir, f"map_{safe_name}.png")
#         try:
#             render_static_png(
#                 site_config, df, map_image_path,
#                 show_numbers=True,
#                 use_mock_fallback=False,  # real data only in reports
#             )
#         except Exception as e:
#             st.warning(f"Could not generate map for {site_address}: {e}")
#             map_image_path = None

#         # ── 3. Open PPTX template ──
#         try:
#             prs = Presentation(template_path)
#         except Exception as e:
#             raise Exception(f"Failed to load PPTX template: {e}")

#         # ── 4. Process each slide ──
#         for slide_idx, slide in enumerate(prs.slides):

#             for shape in slide.shapes:
#                 if not shape.has_text_frame:
#                     continue
#                 for paragraph in shape.text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         txt = run.text

#                         # Slide 0 — Cover letter
#                         if "Name of Resident" in txt:
#                             run.text = txt.replace(
#                                 "Name of Resident",
#                                 f"Resident at {site_address}",
#                             )
#                             txt = run.text
#                         if "Address of Resident" in txt:
#                             run.text = txt.replace(
#                                 "Address of Resident", site_address
#                             )
#                             txt = run.text
#                         if txt.strip() == "Date":
#                             run.text = today_str
#                             txt = run.text

#                         # Slide 3 — Fill [###] PPM
#                         if "[###]" in txt:
#                             low = txt.lower()
#                             if "backyard" in low or "back" in low:
#                                 val = zone_ppm["backyard_ppm"]
#                                 run.text = txt.replace(
#                                     "[###]", str(val) if val else "N/A"
#                                 )
#                             elif "front" in low:
#                                 val = zone_ppm["frontyard_ppm"]
#                                 if val is not None:
#                                     run.text = txt.replace("[###]", str(val))
#                                 else:
#                                     run.text = ""
#                             else:
#                                 run.text = txt.replace(
#                                     "[###]", str(round(site_avg))
#                                 )
#                             txt = run.text

#                         # Legacy placeholders
#                         if "Average Lead concentration (ppm)" in txt:
#                             lbl, _ = get_nysh_category(site_avg)
#                             run.text = txt.replace(
#                                 "Average Lead concentration (ppm)",
#                                 f"Average Lead: {site_avg:.1f} ppm — {lbl}",
#                             )
#                         for old_ph in (
#                             "Visual map of property with color-coded zones",
#                             "Highlight hotspots",
#                             "Heat map of property (no basemap)",
#                             "Heat map of property with basemap",
#                         ):
#                             if old_ph in run.text:
#                                 run.text = run.text.replace(old_ph, "")

#             # --- Map image insertions ---
#             if map_image_path and os.path.exists(map_image_path):
#                 try:
#                     from PIL import Image as PILImage
#                     img = PILImage.open(map_image_path)
#                     img_aspect = img.width / img.height
#                 except Exception:
#                     img_aspect = 1.33  # default 10/7.5

#                 if slide_idx == 2:
#                     max_w, max_h = 6.0, 2.8
#                     if max_w / img_aspect > max_h:
#                         w, h = max_h * img_aspect, max_h
#                     else:
#                         w, h = max_w, max_w / img_aspect
#                     left = Inches((10.0 - w) / 2)
#                     top = Inches(3.8)
#                     try:
#                         slide.shapes.add_picture(
#                             map_image_path, left, top,
#                             width=Inches(w), height=Inches(h),
#                         )
#                     except Exception as e:
#                         st.warning(f"Map insert failed (slide 3) for {site_address}: {e}")

#                 if slide_idx == 3:
#                     max_w, max_h = 7.0, 4.2
#                     if max_w / img_aspect > max_h:
#                         w, h = max_h * img_aspect, max_h
#                     else:
#                         w, h = max_w, max_w / img_aspect
#                     left = Inches((10.0 - w) / 2)
#                     top = Inches(2.6)
#                     try:
#                         slide.shapes.add_picture(
#                             map_image_path, left, top,
#                             width=Inches(w), height=Inches(h),
#                         )
#                     except Exception as e:
#                         st.warning(f"Map insert failed (slide 4) for {site_address}: {e}")

#         # ── 5. Save ──
#         output_file = os.path.join(
#             output_dir, f"Resident_Report_{safe_name}.pptx"
#         )
#         prs.save(output_file)
#         generated_reports.append((site_address, output_file))
#         report_count += 1

#     return report_count, generated_reports


# # ═══════════════════════════════════════════════
# #  UI: HEADER & INSTRUCTIONS
# # ═══════════════════════════════════════════════
# col_header, col_info = st.columns([2, 1])
# with col_header:
#     st.title("⚙️ Data Pipeline Manager")
#     st.markdown("Automated ETL - extraction, transformation, and loading, spatial mapping, and resident report generation.")
# with col_info:
#     st.info(
#         "💡 **Instructions:** Drag and drop your raw XRF `.csv` files below. "
#         "The backend will parse the readings, append them to the Master Database, "
#         "and automatically generate updated site reports."
#     )

# st.markdown("---")

# # ═══════════════════════════════════════════════
# #  UI: UPLOAD & PROCESS
# # ═══════════════════════════════════════════════
# st.subheader("1. Ingest Data")
# uploaded_files = st.file_uploader(
#     "Upload Raw XRF Chemistry Files", type=['csv'], accept_multiple_files=True
# )

# fresh_rebuild = st.checkbox(
#     "🔄 Fresh rebuild (clear old Master Data and reprocess all files)",
#     value=False,
#     help="Check this if you've updated data.py or the XRF Master Data Key "
#          "and need to regenerate everything from scratch.",
# )

# if st.button("🚀 Execute Data Pipeline", type="primary", use_container_width=True):
#     if not uploaded_files:
#         st.warning("⚠️ Please upload at least one chemistry CSV file to begin.")
#     else:
#         with st.status("Executing the Data Pipeline...", expanded=False) as status:
#             try:
#                 # 1. Save uploaded files
#                 xrf_dir = os.path.join("data", "xrf_data")
#                 os.makedirs(xrf_dir, exist_ok=True)
#                 for f in uploaded_files:
#                     with open(os.path.join(xrf_dir, f.name), "wb") as f_out:
#                         f_out.write(f.read())

#                 # 1b. If fresh rebuild, clear old master data
#                 if fresh_rebuild:
#                     master_dir = os.path.join("data", "master_data")
#                     if os.path.exists(master_dir):
#                         import shutil
#                         shutil.rmtree(master_dir)
#                     os.makedirs(master_dir, exist_ok=True)
#                     st.info("🗑️ Cleared old Master Data. Rebuilding from scratch.")

#                 # 2. Run ETL script
#                 result = subprocess.run(
#                     ["python", "src/data.py"], capture_output=True, text=True
#                 )

#                 # Show ETL output for transparency
#                 if result.stdout.strip():
#                     st.code(result.stdout, language="text")

#                 if result.returncode != 0:
#                     status.update(
#                         label="Pipeline Failed during ETL process.", state="error"
#                     )
#                     st.error(f"Backend Error Output:\n{result.stderr}")
#                     st.stop()

#                 # 3. Locate Master Data
#                 master_dir = os.path.join("data", "master_data")
#                 master_files = glob.glob(
#                     os.path.join(master_dir, 'Master_Data_v*.csv')
#                 )
#                 if not master_files:
#                     status.update(
#                         label="Failed to locate output Master Data.", state="error"
#                     )
#                     st.stop()

#                 latest_master = max(
#                     master_files,
#                     key=lambda x: int(
#                         re.search(r'_v(\d+)\.csv', x).group(1)
#                         if re.search(r'_v(\d+)\.csv', x) else 0
#                     ),
#                 )
#                 st.session_state.latest_master_file = latest_master

#                 # 4. Generate Reports
#                 template_path    = os.path.join("src", "Resident_Report_Template.pptx")
#                 site_db_path     = os.path.join(
#                     "data", "site_databases",
#                     "XRF Site Analysis Database W SampleID(Sheet1).csv",
#                 )
#                 site_configs_path = os.path.join(
#                     "data", "site_configs", "site_configs.json"
#                 )
#                 reports_dir  = os.path.join("data", "generated_reports")

#                 if os.path.exists(template_path):
#                     report_count, report_list = generate_pptx_reports(
#                         st.session_state.latest_master_file,
#                         template_path, reports_dir,
#                         site_db_path, site_configs_path,
#                     )
#                     st.session_state.reports_generated = report_count
#                     st.session_state.generated_report_list = report_list
#                 else:
#                     st.warning("⚠️ Template missing. Skipped report generation.")

#                 status.update(
#                     label="Pipeline Execution Complete!", state="complete"
#                 )
#                 st.session_state.pipeline_success = True

#             except Exception as e:
#                 status.update(label="Critical System Error", state="error")
#                 st.error(f"An unexpected error occurred: {str(e)}")
#                 st.session_state.pipeline_success = False

# # ═══════════════════════════════════════════════
# #  UI: RESULTS & EXPORT
# # ═══════════════════════════════════════════════
# if st.session_state.pipeline_success and st.session_state.latest_master_file:
#     st.markdown("---")
#     st.subheader("2. Deployment Artifacts")

#     # KPIs
#     df_result = pd.read_csv(st.session_state.latest_master_file)
#     kpi1, kpi2, kpi3 = st.columns(3)
#     kpi1.metric("Total Records Processed", len(df_result))
#     filled_ids = df_result[
#         df_result['SampleID'].notna() & (df_result['SampleID'] != "")
#     ] if 'SampleID' in df_result.columns else pd.DataFrame()
#     kpi2.metric("Samples Evaluated", filled_ids['SampleID'].nunique() if len(filled_ids) else 0)
#     kpi3.metric("Resident Reports Generated", st.session_state.reports_generated)

#     st.write("")

#     # Download Buttons
#     col_dl1, col_dl2 = st.columns(2)
#     with col_dl1:
#         st.success("### 📊 Master Database")
#         st.caption(
#             f"File: `{os.path.basename(st.session_state.latest_master_file)}`"
#         )
#         with open(st.session_state.latest_master_file, "rb") as file:
#             st.download_button(
#                 label="📥 Download Master Data (CSV)",
#                 data=file,
#                 file_name=os.path.basename(st.session_state.latest_master_file),
#                 mime="text/csv",
#                 use_container_width=True,
#             )

#     with col_dl2:
#         st.info("### 🗂️ Generated Reports")

#         # Gather ALL available reports from the output directory
#         reports_dir = os.path.join("data", "generated_reports")
#         all_report_files = []
#         if os.path.exists(reports_dir):
#             for f in sorted(glob.glob(os.path.join(reports_dir, "Resident_Report_*.pptx")),
#                             key=os.path.getmtime, reverse=True):
#                 fname = os.path.basename(f)
#                 site_name = fname.replace("Resident_Report_", "").replace(".pptx", "")
#                 all_report_files.append((site_name, f))

#         if all_report_files:
#             st.caption(f"{len(all_report_files)} report(s) available.")

#             # Search bar
#             search_query = st.text_input(
#                 "🔍 Search reports by address",
#                 placeholder="e.g. Utica, Cleveland, Schuele...",
#                 key="report_search",
#             )

#             if search_query.strip():
#                 # Show filtered results
#                 filtered = [
#                     (name, path) for name, path in all_report_files
#                     if search_query.strip().lower() in name.lower()
#                 ]

#                 if filtered:
#                     st.caption(f"{len(filtered)} match(es) for \"{search_query}\"")
#                     for idx, (site_name, report_path) in enumerate(filtered):
#                         with open(report_path, "rb") as rpt_file:
#                             st.download_button(
#                                 label=f"📥 {site_name}",
#                                 data=rpt_file,
#                                 file_name=os.path.basename(report_path),
#                                 mime="application/vnd.openxmlformats-officedocument"
#                                      ".presentationml.presentation",
#                                 use_container_width=True,
#                                 key=f"dl_search_{idx}",
#                             )
#                 else:
#                     st.caption(f"No reports matching \"{search_query}\".")

#             else:
#                 # Show only the 4 most recent reports
#                 recent = all_report_files[:4]
#                 st.caption("Showing 4 most recent:")
#                 for idx, (site_name, report_path) in enumerate(recent):
#                     with open(report_path, "rb") as rpt_file:
#                         st.download_button(
#                             label=f"📥 {site_name}",
#                             data=rpt_file,
#                             file_name=os.path.basename(report_path),
#                             mime="application/vnd.openxmlformats-officedocument"
#                                  ".presentationml.presentation",
#                             use_container_width=True,
#                             key=f"dl_recent_{idx}",
#                         )
#         else:
#             st.caption("No reports available yet.")

"""
etl_manager.py — GroundSense Data Pipeline Manager
 
Updates from original:
  1. Imports NYSH colors & rotation helpers from groundsense_config.py
  2. **NEW**: Delegates map image generation to `map_renderer.render_static_png()`
     — the same function site_builder.py uses for its PNG export. PPTX maps
     are now byte-identical to the "No basemap · with numbers" download.
  3. Honors `rotation_deg` saved in site_configs.json (via map_renderer)
  4. Fills in zone-based average Lead PPM values (backyard / front yard)
  5. Data policy: only real XRF data is used in exported maps. Cells without
     matching Master Data render as gray "No Data" — no silent mock fallback.
 
Template slide layout (6 pages, 0-indexed):
  Slide 0: Cover letter — replace "Address of Resident", "Name of Resident", "Date"
  Slide 1: Sample Collection Method — keep as-is
  Slide 2: Soil Report Summary — insert dark map (no basemap)
  Slide 3: Detailed Results — fill [###] with real PPM, insert map
  Slide 4: NY Soil Health info — keep as-is
  Slide 5: Lead level table & safety — keep as-is
"""
 
 
import streamlit as st
import os
import sys
import glob
import re
import subprocess
import json
import math
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from datetime import date
 
# Inject both the script's own directory AND src/ into sys.path so that
# `groundsense_config` and `map_renderer` can be imported regardless of
# whether they live in src/ or at the repo root next to this script.
_HERE = os.path.dirname(os.path.abspath(__file__))
for _p in (_HERE, os.path.join(_HERE, "src")):
    if _p not in sys.path:
        sys.path.insert(0, _p)
 
from groundsense_config import (
    get_nysh_category,
    NYSH_TIERS,
    NYSH_COLORS,
    calculate_coordinate,
    resolve_lod,
)
from map_renderer import render_static_png  # single source of truth for visuals
 
# ═══════════════════════════════════════════════
#  PAGE CONFIGURATION
# ═══════════════════════════════════════════════
st.set_page_config(page_title="ETL Pipeline", page_icon="⚙️", layout="wide")
 
# Initialize Session State
if 'pipeline_success' not in st.session_state:
    st.session_state.pipeline_success = False
if 'latest_master_file' not in st.session_state:
    st.session_state.latest_master_file = None
if 'reports_generated' not in st.session_state:
    st.session_state.reports_generated = 0
if 'generated_report_list' not in st.session_state:
    st.session_state.generated_report_list = []
 
 
# ═══════════════════════════════════════════════
#  ZONE-BASED PPM CALCULATOR
# ═══════════════════════════════════════════════
def compute_zone_averages(site_config, master_df):
    """Compute average Lead PPM for each zone (back, front, yard, transect).
 
    Returns dict, e.g. {"back": 542.3, "front": 718.1}
 
    IMPORTANT: Only uses real XRF data. mock_ppm values are ignored here
    because the resident report must show truthful measurements — placeholders
    belong only in the design-time site_builder preview.
    """
    grid = site_config.get("grid_blocks", {})
 
    if 'LeadPPM_Clean' not in master_df.columns:
        master_df = master_df.copy()
        master_df['LeadPPM_Clean'] = master_df['LeadPPM'].apply(resolve_lod)
 
    zone_values = {}  # zone_name -> [ppm, ...]
 
    for block_id, dims in grid.items():
        if block_id.startswith("_"):
            continue
        zone = dims.get("zone", "yard")
        patterns = dims.get("sample_id_patterns", [])
 
        ppm = None
        for pat in patterns:
            if not pat:
                continue
            matches = master_df[
                master_df['SampleID'].str.contains(pat, case=False, na=False)
            ]
            if not matches.empty:
                avg = matches['LeadPPM_Clean'].mean()
                if pd.notna(avg):
                    ppm = avg
                    break
 
        # Skip mock fallback — real data only
        if ppm is not None and not (isinstance(ppm, float) and math.isnan(ppm)):
            zone_values.setdefault(zone, []).append(ppm)
 
    return {z: sum(v) / len(v) for z, v in zone_values.items() if v}
 
 
def format_zone_ppm(zone_averages):
    """Map zone averages to backyard_ppm / frontyard_ppm for template filling.
 
    Single-zone sites (yard, transect) map to backyard only.
    """
    result = {"backyard_ppm": None, "frontyard_ppm": None}
    for zone, avg in zone_averages.items():
        z = zone.lower()
        if z in ("back", "backyard"):
            result["backyard_ppm"] = round(avg)
        elif z in ("front", "frontyard", "front_yard"):
            result["frontyard_ppm"] = round(avg)
        elif z in ("yard", "transect"):
            result["backyard_ppm"] = round(avg)
    return result
 
 
# ═══════════════════════════════════════════════
#  REPORT GENERATION (PPTX)
# ═══════════════════════════════════════════════
def generate_pptx_reports(master_csv_path, template_path, output_dir,
                          site_configs_path):
    """Generate one PPTX resident report per site address.
 
    Uses the shared `map_renderer.render_static_png()` so the map image
    embedded in the PPTX is byte-identical to the PNG download from
    site_builder.py's "No basemap · with numbers" option.
 
    Returns (report_count, list_of (site_address, filepath) tuples).
    """
    df = pd.read_csv(master_csv_path)
    df = df[df['SampleID'].notna() & (df['SampleID'] != "")]
    df['LeadPPM_Clean'] = df['LeadPPM'].apply(resolve_lod)
 
    # Load site configs — one report per site address
    site_configs = {}
    if os.path.exists(site_configs_path):
        with open(site_configs_path, 'r') as f:
            raw = json.load(f)
        site_configs = {s["address"]: s for s in raw}
 
    if not site_configs:
        st.warning("⚠️ No site configurations found. Cannot generate reports.")
        return 0, []
 
    os.makedirs(output_dir, exist_ok=True)
    maps_dir = os.path.join(output_dir, "map_images")
    os.makedirs(maps_dir, exist_ok=True)
 
    report_count = 0
    generated_reports = []
    today_str = date.today().strftime("%m/%d/%Y")
 
    for site_address, site_config in site_configs.items():
 
        # ── 1. Compute zone averages (real data only) ──
        zone_averages = compute_zone_averages(site_config, df)
        zone_ppm = format_zone_ppm(zone_averages)
 
        all_zone_vals = list(zone_averages.values())
        if all_zone_vals:
            site_avg = sum(all_zone_vals) / len(all_zone_vals)
        else:
            # No real data for this site — skip report generation entirely.
            # Previously we fell back to mock_ppm, but that leaks placeholder
            # numbers into official reports, which we explicitly don't want.
            st.info(
                f"⏭️  Skipped '{site_address}' — no real XRF data matched yet."
            )
            continue
 
        if zone_ppm["backyard_ppm"] is None:
            zone_ppm["backyard_ppm"] = round(site_avg)
 
        # ── 2. Render static map via shared renderer ──
        safe_name = "".join(
            c for c in site_address if c.isalnum() or c == ' '
        ).rstrip()
        map_image_path = os.path.join(maps_dir, f"map_{safe_name}.png")
        try:
            render_static_png(
                site_config, df, map_image_path,
                show_numbers=True,
                use_mock_fallback=False,  # real data only in reports
            )
        except Exception as e:
            st.warning(f"Could not generate map for {site_address}: {e}")
            map_image_path = None
 
        # ── 3. Open PPTX template ──
        try:
            prs = Presentation(template_path)
        except Exception as e:
            raise Exception(f"Failed to load PPTX template: {e}")
 
        # ── 4. Process each slide ──
        for slide_idx, slide in enumerate(prs.slides):
 
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        txt = run.text
 
                        # Slide 0 — Cover letter
                        if "Name of Resident" in txt:
                            run.text = txt.replace(
                                "Name of Resident",
                                f"Resident at {site_address}",
                            )
                            txt = run.text
                        if "Address of Resident" in txt:
                            run.text = txt.replace(
                                "Address of Resident", site_address
                            )
                            txt = run.text
                        if txt.strip() == "Date":
                            run.text = today_str
                            txt = run.text
 
                        # Slide 3 — Fill [###] PPM
                        if "[###]" in txt:
                            low = txt.lower()
                            if "backyard" in low or "back" in low:
                                val = zone_ppm["backyard_ppm"]
                                run.text = txt.replace(
                                    "[###]", str(val) if val else "N/A"
                                )
                            elif "front" in low:
                                val = zone_ppm["frontyard_ppm"]
                                if val is not None:
                                    run.text = txt.replace("[###]", str(val))
                                else:
                                    run.text = ""
                            else:
                                run.text = txt.replace(
                                    "[###]", str(round(site_avg))
                                )
                            txt = run.text
 
                        # Legacy placeholders
                        if "Average Lead concentration (ppm)" in txt:
                            lbl, _ = get_nysh_category(site_avg)
                            run.text = txt.replace(
                                "Average Lead concentration (ppm)",
                                f"Average Lead: {site_avg:.1f} ppm — {lbl}",
                            )
                        for old_ph in (
                            "Visual map of property with color-coded zones",
                            "Highlight hotspots",
                            "Heat map of property (no basemap)",
                            "Heat map of property with basemap",
                        ):
                            if old_ph in run.text:
                                run.text = run.text.replace(old_ph, "")
 
            # --- Map image insertions ---
            if map_image_path and os.path.exists(map_image_path):
                try:
                    from PIL import Image as PILImage
                    img = PILImage.open(map_image_path)
                    img_aspect = img.width / img.height
                except Exception:
                    img_aspect = 1.33  # default 10/7.5
 
                if slide_idx == 2:
                    max_w, max_h = 6.0, 2.8
                    if max_w / img_aspect > max_h:
                        w, h = max_h * img_aspect, max_h
                    else:
                        w, h = max_w, max_w / img_aspect
                    left = Inches((10.0 - w) / 2)
                    top = Inches(3.8)
                    try:
                        slide.shapes.add_picture(
                            map_image_path, left, top,
                            width=Inches(w), height=Inches(h),
                        )
                    except Exception as e:
                        st.warning(f"Map insert failed (slide 3) for {site_address}: {e}")
 
                if slide_idx == 3:
                    max_w, max_h = 7.0, 4.2
                    if max_w / img_aspect > max_h:
                        w, h = max_h * img_aspect, max_h
                    else:
                        w, h = max_w, max_w / img_aspect
                    left = Inches((10.0 - w) / 2)
                    top = Inches(2.6)
                    try:
                        slide.shapes.add_picture(
                            map_image_path, left, top,
                            width=Inches(w), height=Inches(h),
                        )
                    except Exception as e:
                        st.warning(f"Map insert failed (slide 4) for {site_address}: {e}")
 
        # ── 5. Save ──
        output_file = os.path.join(
            output_dir, f"Resident_Report_{safe_name}.pptx"
        )
        prs.save(output_file)
        generated_reports.append((site_address, output_file))
        report_count += 1
 
    return report_count, generated_reports
 
 
# ═══════════════════════════════════════════════
#  UI: HEADER & INSTRUCTIONS
# ═══════════════════════════════════════════════
col_header, col_info = st.columns([2, 1])
with col_header:
    st.title("⚙️ Data Pipeline Manager")
    st.markdown("Automated ETL (Extraction, Transformation, and Loading), spatial mapping, and resident report generation.")
with col_info:
    st.info(
        "💡 **Instructions:** Drag and drop your raw XRF `.csv` files below. "
        "The backend will parse the readings, append them to the Master Database, "
        "and automatically generate updated site reports."
    )
 
st.markdown("---")
 
# ═══════════════════════════════════════════════
#  UI: UPLOAD & PROCESS
# ═══════════════════════════════════════════════
st.subheader("1. Ingest Data")
uploaded_files = st.file_uploader(
    "Upload Raw XRF Chemistry Files", type=['csv'], accept_multiple_files=True
)
 
fresh_rebuild = st.checkbox(
    "🔄 Fresh rebuild (clear old Master Data and reprocess all files)",
    value=False,
    help="Check this if you've updated data.py or the XRF Master Data Key "
         "and need to regenerate everything from scratch.",
)
 
if st.button("🚀 Execute Data Pipeline", type="primary", use_container_width=True):
    if not uploaded_files:
        st.warning("⚠️ Please upload at least one chemistry CSV file to begin.")
    else:
        with st.status("Executing the Data Pipeline...", expanded=False) as status:
            try:
                # 1. Save uploaded files
                xrf_dir = os.path.join("data", "xrf_data")
                os.makedirs(xrf_dir, exist_ok=True)
                for f in uploaded_files:
                    with open(os.path.join(xrf_dir, f.name), "wb") as f_out:
                        f_out.write(f.read())
 
                # 1b. If fresh rebuild, clear old master data
                if fresh_rebuild:
                    master_dir = os.path.join("data", "master_data")
                    if os.path.exists(master_dir):
                        import shutil
                        shutil.rmtree(master_dir)
                    os.makedirs(master_dir, exist_ok=True)
                    st.info("🗑️ Cleared old Master Data. Rebuilding from scratch.")
 
                # 2. Run ETL script
                result = subprocess.run(
                    ["python", "src/data.py"], capture_output=True, text=True
                )
 
                # Show ETL output for transparency
                if result.stdout.strip():
                    st.code(result.stdout, language="text")
 
                if result.returncode != 0:
                    status.update(
                        label="Pipeline Failed during ETL process.", state="error"
                    )
                    st.error(f"Backend Error Output:\n{result.stderr}")
                    st.stop()
 
                # 3. Locate Master Data
                master_dir = os.path.join("data", "master_data")
                master_files = glob.glob(
                    os.path.join(master_dir, 'Master_Data_v*.csv')
                )
                if not master_files:
                    status.update(
                        label="Failed to locate output Master Data.", state="error"
                    )
                    st.stop()
 
                latest_master = max(
                    master_files,
                    key=lambda x: int(
                        re.search(r'_v(\d+)\.csv', x).group(1)
                        if re.search(r'_v(\d+)\.csv', x) else 0
                    ),
                )
                st.session_state.latest_master_file = latest_master
 
                # 4. Generate Reports
                template_path    = os.path.join("src", "Resident_Report_Template.pptx")
                # NOTE: site_db_path was previously passed here pointing at
                # `XRF Site Analysis Database W SampleID(Sheet1).csv`, but
                # generate_pptx_reports never read from it — addresses come
                # from site_configs.json. Removed to avoid re-introducing a
                # dependency on the (PII-scrubbed) field CSV.
                site_configs_path = os.path.join(
                    "data", "site_configs", "site_configs.json"
                )
                reports_dir  = os.path.join("data", "generated_reports")
 
                if os.path.exists(template_path):
                    report_count, report_list = generate_pptx_reports(
                        st.session_state.latest_master_file,
                        template_path, reports_dir,
                        site_configs_path,
                    )
                    st.session_state.reports_generated = report_count
                    st.session_state.generated_report_list = report_list
                else:
                    st.warning("⚠️ Template missing. Skipped report generation.")
 
                status.update(
                    label="Pipeline Execution Complete!", state="complete"
                )
                st.session_state.pipeline_success = True
 
            except Exception as e:
                status.update(label="Critical System Error", state="error")
                st.error(f"An unexpected error occurred: {str(e)}")
                st.session_state.pipeline_success = False
 
# ═══════════════════════════════════════════════
#  UI: RESULTS & EXPORT
# ═══════════════════════════════════════════════
if st.session_state.pipeline_success and st.session_state.latest_master_file:
    st.markdown("---")
    st.subheader("2. Deployment Artifacts")
 
    # KPIs
    df_result = pd.read_csv(st.session_state.latest_master_file)
    kpi1, kpi2, kpi3 = st.columns(3)
    kpi1.metric("Total Records Processed", len(df_result))
    filled_ids = df_result[
        df_result['SampleID'].notna() & (df_result['SampleID'] != "")
    ] if 'SampleID' in df_result.columns else pd.DataFrame()
    kpi2.metric("Samples Evaluated", filled_ids['SampleID'].nunique() if len(filled_ids) else 0)
    kpi3.metric("Resident Reports Generated", st.session_state.reports_generated)
 
    st.write("")
 
    # Download Buttons
    col_dl1, col_dl2 = st.columns(2)
    with col_dl1:
        st.success("### 📊 Master Database")
        st.caption(
            f"File: `{os.path.basename(st.session_state.latest_master_file)}`"
        )
        with open(st.session_state.latest_master_file, "rb") as file:
            st.download_button(
                label="📥 Download Master Data (CSV)",
                data=file,
                file_name=os.path.basename(st.session_state.latest_master_file),
                mime="text/csv",
                use_container_width=True,
            )
 
    with col_dl2:
        st.info("### 🗂️ Generated Reports")
 
        # Gather ALL available reports from the output directory
        reports_dir = os.path.join("data", "generated_reports")
        all_report_files = []
        if os.path.exists(reports_dir):
            for f in sorted(glob.glob(os.path.join(reports_dir, "Resident_Report_*.pptx")),
                            key=os.path.getmtime, reverse=True):
                fname = os.path.basename(f)
                site_name = fname.replace("Resident_Report_", "").replace(".pptx", "")
                all_report_files.append((site_name, f))
 
        if all_report_files:
            st.caption(f"{len(all_report_files)} report(s) available.")
 
            # Search bar
            search_query = st.text_input(
                "🔍 Search reports by address",
                placeholder="e.g. Utica, Cleveland, Schuele...",
                key="report_search",
            )
 
            if search_query.strip():
                # Show filtered results
                filtered = [
                    (name, path) for name, path in all_report_files
                    if search_query.strip().lower() in name.lower()
                ]
 
                if filtered:
                    st.caption(f"{len(filtered)} match(es) for \"{search_query}\"")
                    for idx, (site_name, report_path) in enumerate(filtered):
                        with open(report_path, "rb") as rpt_file:
                            st.download_button(
                                label=f"📥 {site_name}",
                                data=rpt_file,
                                file_name=os.path.basename(report_path),
                                mime="application/vnd.openxmlformats-officedocument"
                                     ".presentationml.presentation",
                                use_container_width=True,
                                key=f"dl_search_{idx}",
                            )
                else:
                    st.caption(f"No reports matching \"{search_query}\".")
 
            else:
                # Show only the 4 most recent reports
                recent = all_report_files[:4]
                st.caption("Showing 4 most recent:")
                for idx, (site_name, report_path) in enumerate(recent):
                    with open(report_path, "rb") as rpt_file:
                        st.download_button(
                            label=f"📥 {site_name}",
                            data=rpt_file,
                            file_name=os.path.basename(report_path),
                            mime="application/vnd.openxmlformats-officedocument"
                                 ".presentationml.presentation",
                            use_container_width=True,
                            key=f"dl_recent_{idx}",
                        )
        else:
            st.caption("No reports available yet.")